import json as json_module
import os
import tempfile
import time
from datetime import datetime, timezone
from pathlib import Path

from fastapi import APIRouter, Depends, HTTPException, Query, UploadFile, File
from openai import OpenAI
from sqlalchemy.orm import Session

from .. import auth as auth_module
from .. import crud, schemas
from ..config import (
    IMPORT_CHUNK_SIZE,
    IMPORT_MAX_CHUNKS,
    IMPORT_UPSTREAM_TIMEOUT,
    OPENAI_API_KEY,
    OPENAI_BASE_URL,
    OPENAI_MODEL,
)
from ..crud import derive_course_name_from_filename
from ..database import get_db
from ..models import Question as QuestionModel
from ..schemas import ImportedQuestion
from ..utils import normalize_answer, VALID_QUESTION_TYPES

router = APIRouter(prefix="/imports", tags=["imports"])

ALLOWED_EXTENSIONS = {".docx", ".pptx"}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB
AI_CHUNK_SIZE = max(1000, IMPORT_CHUNK_SIZE)  # characters per chunk
MAX_CHUNKS = max(1, IMPORT_MAX_CHUNKS)  # maximum chunks to avoid runaway costs


def _elapsed_ms(start: float) -> int:
    return max(0, round((time.perf_counter() - start) * 1000))


def _build_timing(
    *,
    total_start: float,
    extract_ms: int = 0,
    parse_timing: dict | None = None,
) -> schemas.ImportTiming:
    parse_timing = parse_timing or {}
    return schemas.ImportTiming(
        extract_ms=extract_ms,
        chunk_ms=int(parse_timing.get("chunk_ms") or 0),
        ai_ms=int(parse_timing.get("ai_ms") or 0),
        total_ms=_elapsed_ms(total_start),
        chunks=int(parse_timing.get("chunks") or 0),
        ai_chunks=list(parse_timing.get("ai_chunks") or []),
    )


# ── Text extraction (enhanced) ──────────────────────────────────────────────

def extract_text_and_warnings(file_path: str) -> tuple[str, list[str]]:
    """Extract text from docx/pptx, also extracting tables and detecting images.
    Returns (text, warnings).
    """
    ext = Path(file_path).suffix.lower()
    warnings: list[str] = []

    if ext == ".docx":
        return _extract_docx(file_path, warnings)
    elif ext == ".pptx":
        return _extract_pptx(file_path, warnings)
    raise ValueError(f"不支持的文件格式: {ext}")


def _extract_docx(path: str, warnings: list[str]) -> tuple[str, list[str]]:
    from docx import Document
    doc = Document(path)
    parts: list[str] = []

    # Paragraphs
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append(t)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                parts.append(" | ".join(row_texts))

    # Detect images
    has_image = False
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            has_image = True
            break
    if has_image:
        warnings.append("文档中包含图片，暂不支持 OCR 识别图片中的文字，请手动核对")

    return "\n".join(parts), warnings


def _extract_pptx(path: str, warnings: list[str]) -> tuple[str, list[str]]:
    from pptx import Presentation
    prs = Presentation(path)
    parts: list[str] = []
    has_image = False

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
            if shape.shape_type is not None:
                # check for image-like shapes
                try:
                    if "Picture" in type(shape).__name__ or hasattr(shape, "image"):
                        has_image = True
                except Exception:
                    pass
            # Table in PPTX
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        parts.append(" | ".join(row_texts))

    if has_image:
        warnings.append("演示文稿中包含图片，暂不支持 OCR 识别图片中的文字，请手动核对")

    return "\n".join(parts), warnings


# Backward-compatible wrapper
def extract_text_from_file(file_path: str) -> str:
    text, _ = extract_text_and_warnings(file_path)
    return text


# ── AI parsing helpers ──────────────────────────────────────────────────────

def _validate_question_item(item: dict) -> tuple[dict | None, str | None]:
    """Validate a single parsed question dict. Returns (normalized_item, error_message)."""
    q_type = (item.get("type") or "").strip()
    question = (item.get("question") or "").strip()
    answer = (item.get("answer") or "").strip()

    if q_type not in VALID_QUESTION_TYPES:
        return None, f"无效的题目类型 '{q_type}'"
    if not question:
        return None, "题目题干不能为空"
    if not answer:
        return None, "题目答案不能为空"
    if q_type in ("single_choice", "multiple_choice"):
        opts = item.get("options")
        if not opts or not isinstance(opts, dict) or len(opts) < 2:
            return None, f"选择题（{q_type}）必须至少提供两个选项"

    return {
        "type": q_type,
        "question": question,
        "options": item.get("options"),
        "answer": normalize_answer(answer, q_type),
        "analysis": (item.get("analysis") or "").strip(),
        "subject": (item.get("subject") or "默认科目").strip(),
        "chapter": (item.get("chapter") or "默认章节").strip(),
        "difficulty": (item.get("difficulty") or "normal").strip(),
    }, None


def _extract_questions_from_ai_response(raw: str) -> tuple[list[dict], list[str]]:
    """Extract a list of question dicts from an AI JSON response.
    Returns (questions, warnings).
    """
    warnings: list[str] = []
    try:
        parsed = json_module.loads(raw)
    except json_module.JSONDecodeError:
        warnings.append("AI 返回了非 JSON 格式内容，已忽略此分块")
        return [], warnings

    if isinstance(parsed, list):
        items = parsed
    elif isinstance(parsed, dict):
        items = (
            parsed.get("questions")
            or parsed.get("items")
            or parsed.get("data")
            or parsed.get("result")
            or []
        )
        if not items and parsed:
            if "question" in parsed or "type" in parsed:
                items = [parsed]
    else:
        items = []

    return items, warnings


def _call_ai_parse_chunk(text_chunk: str, chunk_index: int) -> tuple[list[dict], list[str]]:
    """Call AI to parse a single text chunk. Returns (question_dicts, warnings)."""
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url=OPENAI_BASE_URL,
        timeout=IMPORT_UPSTREAM_TIMEOUT,
    )
    prompt = (
        "你是一个题目解析助手。请将以下教育文档内容解析为 JSON 数组，每道题为一个对象。\n\n"
        "支持的题目类型：\n"
        "- single_choice（单选题）\n"
        "- multiple_choice（多选题）\n"
        "- true_false（判断题）\n"
        "- fill_blank（填空题）\n"
        "- short_answer（简答题）\n\n"
        "每个对象字段：\n"
        "- type: 题目类型（必填）\n"
        "- question: 题目题干（必填）\n"
        "- options: 选择题的选项字典，如 {\"A\": \"内容1\", \"B\": \"内容2\"}（选择题必填）\n"
        "- answer: 正确答案（必填）\n"
        "- analysis: 解析（可选）\n"
        "- subject: 科目（可选，默认\"默认科目\"）\n"
        "- chapter: 章节（可选，默认\"默认章节\"）\n"
        "- difficulty: 难度 easy/normal/hard（可选，默认\"normal\"）\n\n"
        "请以 JSON 格式输出，包含一个 questions 字段，其值为题目数组。\n"
        '示例：{"questions": [{"type": "single_choice", "question": "...", "options": {"A": "...", "B": "..."}, "answer": "A"}]}\n\n'
        f"文档内容：\n{text_chunk}"
    )

    # Override the legacy garbled prompt with a clean prompt. Keep it mostly
    # ASCII so future Windows encoding issues are less likely to break AI import.
    prompt = (
        "You are an exam-question extraction assistant. Convert the document text into strict JSON.\n"
        "Return ONLY a JSON object with this shape:\n"
        "{\n"
        '  "questions": [\n'
        "    {\n"
        '      "type": "single_choice | multiple_choice | true_false | fill_blank | short_answer",\n'
        '      "question": "question text",\n'
        '      "options": {"A": "option A", "B": "option B"},\n'
        '      "answer": "correct answer",\n'
        '      "analysis": "short explanation",\n'
        '      "subject": "subject name",\n'
        '      "chapter": "chapter name",\n'
        '      "difficulty": "easy | normal | hard"\n'
        "    }\n"
        "  ]\n"
        "}\n\n"
        "Rules:\n"
        "1. Use the original language of the document for question text and analysis.\n"
        "2. For choice questions, options must be an object keyed by A/B/C/D.\n"
        "3. For true_false answers, use one of: true, false, yes, no.\n"
        "4. Do not include markdown fences or explanations outside JSON.\n\n"
        f"Document text:\n{text_chunk}"
    )

    response = client.chat.completions.create(
        model=OPENAI_MODEL,
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"},
        temperature=0.1,
        max_tokens=3000,
    )
    raw = response.choices[0].message.content
    if not raw:
        return [], [f"第 {chunk_index+1} 部分 AI 返回了空响应"]

    items, warn = _extract_questions_from_ai_response(raw)
    if warn:
        warn = [f"第 {chunk_index+1} 部分: {w}" for w in warn]
    return items, warn


def _deduplicate_questions(questions: list[dict]) -> list[dict]:
    """Deduplicate questions by (normalized) question text."""
    seen: set[str] = set()
    result: list[dict] = []
    for q in questions:
        key = (q.get("question") or "").strip()[:100]
        if key and key not in seen:
            seen.add(key)
            result.append(q)
    return result


def call_ai_parse(
    text: str,
) -> tuple[list[dict], list[str], dict]:
    """Parse document text into question dicts with chunking + dedup.

    Returns (all_valid_questions, warnings, timing).
    """
    total_start = time.perf_counter()
    timing = {
        "chunk_ms": 0,
        "ai_ms": 0,
        "total_ms": 0,
        "chunks": 0,
        "ai_chunks": [],
    }
    if not OPENAI_API_KEY:
        raise HTTPException(
            status_code=400,
            detail="未配置 OPENAI_API_KEY，请在 .env 文件中设置 OPENAI_API_KEY 以使用 AI 自动导入功能",
        )

    all_items: list[dict] = []
    all_warnings: list[str] = []

    # Chunk the text
    chunk_start = time.perf_counter()
    paragraphs = [p for p in text.split("\n") if p.strip()]
    chunks: list[str] = []
    current = ""
    for para in paragraphs:
        if len(current) + len(para) + 1 > AI_CHUNK_SIZE:
            if current:
                chunks.append(current)
            current = para
        else:
            current = current + "\n" + para if current else para
    if current:
        chunks.append(current)
    timing["chunk_ms"] = _elapsed_ms(chunk_start)

    # Limit chunks
    if len(chunks) > MAX_CHUNKS:
        all_warnings.append(f"文档过长，仅处理前 {MAX_CHUNKS} 部分（共 {len(chunks)} 部分）")
        chunks = chunks[:MAX_CHUNKS]
    timing["chunks"] = len(chunks)

    # Parse each chunk
    for i, chunk in enumerate(chunks):
        ai_start = time.perf_counter()
        try:
            items, warns = _call_ai_parse_chunk(chunk, i)
            all_items.extend(items)
            all_warnings.extend(warns)
        except Exception as e:
            all_warnings.append(f"第 {i+1} 部分解析失败: {e}")
        finally:
            ai_chunk_ms = _elapsed_ms(ai_start)
            timing["ai_chunks"].append(ai_chunk_ms)
            timing["ai_ms"] += ai_chunk_ms

    # Deduplicate
    all_items = _deduplicate_questions(all_items)

    # Validate each item
    valid: list[dict] = []
    for idx, item in enumerate(all_items):
        validated, err = _validate_question_item(item)
        if validated:
            validated["line_number"] = idx + 1
            valid.append(validated)
        else:
            all_warnings.append(f"第 {idx+1} 题格式有误: {err}")

    timing["total_ms"] = _elapsed_ms(total_start)
    return valid, all_warnings, timing


# ── POST /file (unchanged) ────────────────────────────────────────────────

@router.post("/file", response_model=schemas.FileExtractResponse)
async def upload_file(
    file: UploadFile = File(...),
    course_id: int = Query(0, ge=0, description="目标课程ID（0表示自动创建/使用「未分类题库」）"),
    db: Session = Depends(get_db),
    _current_user=Depends(auth_module.get_current_user),
):
    total_start = time.perf_counter()
    ext = Path(file.filename or "").suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件格式 '{ext}'，仅支持 .docx 或 .pptx 文件",
        )

    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"文件大小超过限制（最大 {MAX_FILE_SIZE // (1024*1024)}MB）",
        )

    tmp_path = None
    try:
        suffix = ext or ".tmp"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        extract_start = time.perf_counter()
        text, _ = extract_text_and_warnings(tmp_path)
        extract_ms = _elapsed_ms(extract_start)
        return schemas.FileExtractResponse(
            text=text,
            filename=file.filename or "unknown",
            suggested_course_name=derive_course_name_from_filename(file.filename or ""),
            timing=_build_timing(total_start=total_start, extract_ms=extract_ms),
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文件提取失败: {str(e)}")
    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


# ── POST /file/preview (parse only, no DB writes) ──────────────────────────

@router.post("/file/preview", response_model=schemas.PreviewImportResponse)
async def preview_import(
    file: UploadFile = File(...),
    _current_user=Depends(auth_module.get_current_user),
):
    """Upload a docx/pptx, extract text, parse via AI, return preview.
    Does NOT write anything to the database.
    """
    total_start = time.perf_counter()
    ext = Path(file.filename or "").suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件格式 '{ext}'，仅支持 .docx 或 .pptx 文件",
        )

    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"文件大小超过限制（最大 {MAX_FILE_SIZE // (1024*1024)}MB）",
        )

    tmp_path = None
    try:
        suffix = ext or ".tmp"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        extract_start = time.perf_counter()
        text, extract_warnings = extract_text_and_warnings(tmp_path)
        extract_ms = _elapsed_ms(extract_start)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文件提取失败: {str(e)}")
    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    if not text.strip():
        raise HTTPException(status_code=400, detail="文件中未提取到任何文本内容")

    # AI parse
    questions, ai_warnings, parse_timing = call_ai_parse(text)
    all_warnings = extract_warnings + ai_warnings
    suggested = derive_course_name_from_filename(file.filename or "")

    total_parsed = len(questions)
    valid = questions
    total_valid = len(valid)
    total_invalid = 0
    # Count invalids from warnings (rough: filter warnings matching "格式有误")
    for w in ai_warnings:
        if "格式有误" in w:
            total_invalid += 1

    return schemas.PreviewImportResponse(
        questions=[ImportedQuestion(**q) for q in valid],
        suggested_course_name=suggested,
        warnings=all_warnings,
        total_parsed=total_parsed,
        total_valid=total_valid,
        total_invalid=total_invalid,
        timing=_build_timing(
            total_start=total_start,
            extract_ms=extract_ms,
            parse_timing=parse_timing,
        ),
    )


# ── POST /imports/confirm (single-transaction import) ─────────────────────

@router.post("/confirm", response_model=schemas.ConfirmImportResponse)
def confirm_import(
    body: schemas.ConfirmImportRequest,
    db: Session = Depends(get_db),
    current_user=Depends(auth_module.get_current_user),
):
    """Confirm and persist user-edited questions from preview.

    All questions are validated first. If any fail, nothing is written.
    Course resolution happens only after all questions pass.
    """
    if not body.questions:
        raise HTTPException(status_code=422, detail="没有待导入的题目")

    # 1. Validate all questions before touching the DB
    errors: list[str] = []
    validated_items: list[schemas.ImportedQuestion] = []
    for i, q in enumerate(body.questions):
        d = q.model_dump()
        validated, err = _validate_question_item(d)
        if err:
            errors.append(f"第 {i+1} 题: {err}")
        else:
            validated_items.append(schemas.ImportedQuestion(**validated))

    if errors:
        raise HTTPException(
            status_code=422,
            detail=f"部分题目校验不通过，未导入任何题目：{'；'.join(errors)}",
        )

    # 2. Resolve target course (only after all questions pass)
    try:
        bank, _ = crud.resolve_course(
            db, current_user.id,
            course_id=body.course_id,
            course_name=body.course_name.strip() if body.course_name else None,
        )
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))

    # 3. Single-transaction import
    imported = 0
    try:
        for q in validated_items:
            question = QuestionModel(
                owner_id=current_user.id,
                course_id=bank.id,
                visibility="private",
                source="import",
                created_at=datetime.now(timezone.utc),
                subject=q.subject,
                chapter=q.chapter,
                type=q.type,
                question=q.question,
                answer=normalize_answer(q.answer, q.type),
                analysis=q.analysis or "",
                difficulty=q.difficulty or "normal",
            )
            question.set_options_dict(q.options)
            db.add(question)
            imported += 1
        db.commit()
    except Exception as e:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"导入失败，已回滚: {e}")

    return schemas.ConfirmImportResponse(
        imported_count=imported,
        course_id=bank.id,
        course_name=bank.name,
    )


# ── POST /file/auto (refactored to reuse shared logic) ─────────────────────

@router.post("/file/auto")
async def import_file_auto(
    file: UploadFile = File(...),
    course_id: int = Query(0, ge=0, description="目标课程ID（0表示使用 course_name 或文件名推断）"),
    course_name: str = Query("", description="目标课程名（course_id=0 时优先使用；空则从文件名推断）"),
    db: Session = Depends(get_db),
    current_user=Depends(auth_module.get_current_user),
):
    total_start = time.perf_counter()
    ext = Path(file.filename or "").suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件格式 '{ext}'，仅支持 .docx 或 .pptx 文件",
        )

    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"文件大小超过限制（最大 {MAX_FILE_SIZE // (1024*1024)}MB）",
        )

    tmp_path = None
    try:
        suffix = ext or ".tmp"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        extract_start = time.perf_counter()
        text, _ = extract_text_and_warnings(tmp_path)
        extract_ms = _elapsed_ms(extract_start)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文件提取失败: {str(e)}")
    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    if not text.strip():
        raise HTTPException(status_code=400, detail="文件中未提取到任何文本内容")

    # AI parse first (no DB writes yet)
    question_dicts, _warnings, parse_timing = call_ai_parse(text)
    if not question_dicts:
        detail = "未能从文档中解析出任何有效题目"
        if _warnings:
            detail += "：" + "；".join(_warnings[:3])
        raise HTTPException(status_code=400, detail=detail)

    # Resolve target course only after successful parsing
    try:
        if course_id > 0:
            bank, _ = crud.resolve_course(db, current_user.id, course_id=course_id)
        elif course_name.strip():
            bank, _ = crud.resolve_course(db, current_user.id, course_name=course_name.strip())
        else:
            derived = derive_course_name_from_filename(file.filename or "")
            bank, _ = crud.resolve_course(db, current_user.id, course_name=derived)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))

    # Import — single transaction
    imported = 0
    for d in question_dicts:
        try:
            question = QuestionModel(
                owner_id=current_user.id,
                course_id=bank.id,
                visibility="private",
                source="import",
                created_at=datetime.now(timezone.utc),
                subject=d["subject"],
                chapter=d["chapter"],
                type=d["type"],
                question=d["question"],
                answer=normalize_answer(d["answer"], d["type"]),
                analysis=d.get("analysis", ""),
                difficulty=d.get("difficulty", "normal"),
            )
            question.set_options_dict(d.get("options"))
            db.add(question)
            imported += 1
        except Exception as exc:
            _warnings.append(f"导入异常: {exc}")

    if imported > 0:
        db.commit()

    return schemas.FileAutoResponse(
        imported_count=imported,
        course_id=bank.id,
        course_name=bank.name,
        timing=_build_timing(
            total_start=total_start,
            extract_ms=extract_ms,
            parse_timing=parse_timing,
        ),
    )
