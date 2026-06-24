import json as json_module
import os
import re
import tempfile
import time
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from fastapi import HTTPException, UploadFile
from openai import OpenAI
from sqlalchemy.orm import Session

from .. import crud, schemas
from ..config import (
    IMPORT_CHUNK_SIZE,
    IMPORT_MAX_CHUNKS,
    IMPORT_UPSTREAM_TIMEOUT,
    OPENAI_API_KEY,
    OPENAI_BASE_URL,
    OPENAI_MODEL,
)
from ..crud import derive_course_name_from_filename
from ..models import Question as QuestionModel
from ..utils import VALID_QUESTION_TYPES, normalize_answer

try:
    from openai import APITimeoutError
except ImportError:  # pragma: no cover
    APITimeoutError = TimeoutError


ALLOWED_EXTENSIONS = {".docx", ".pptx"}
MAX_FILE_SIZE = 10 * 1024 * 1024
AI_CHUNK_SIZE = max(1000, IMPORT_CHUNK_SIZE)
MAX_CHUNKS = max(1, IMPORT_MAX_CHUNKS)


@dataclass
class SavedUpload:
    filename: str
    ext: str
    path: str


def elapsed_ms(start: float) -> int:
    return max(0, round((time.perf_counter() - start) * 1000))


def build_timing(
    *,
    total_start: float,
    extract_ms: int = 0,
    parse_timing: dict | None = None,
) -> schemas.ImportTiming:
    parse_timing = parse_timing or {}
    return schemas.ImportTiming(
        extract_ms=extract_ms,
        chunk_ms=int(parse_timing.get("chunk_ms") or 0),
        ai_ms=int(parse_timing.get("ai_ms") or 0),
        total_ms=elapsed_ms(total_start),
        chunks=int(parse_timing.get("chunks") or 0),
        ai_chunks=list(parse_timing.get("ai_chunks") or []),
    )


def sync_ai_settings(*, api_key=None, base_url=None, client_class=None) -> None:
    global OPENAI_API_KEY, OPENAI_BASE_URL, OpenAI
    if api_key is not None:
        OPENAI_API_KEY = api_key
    if base_url is not None:
        OPENAI_BASE_URL = base_url
    if client_class is not None:
        OpenAI = client_class


def validate_upload(filename: str, content: bytes) -> str:
    ext = Path(filename or "").suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件格式 '{ext}'，仅支持 .docx 或 .pptx 文件",
        )
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"文件大小超过限制（最大 {MAX_FILE_SIZE // (1024 * 1024)}MB）",
        )
    return ext


async def save_upload_to_temp(file: UploadFile) -> SavedUpload:
    filename = file.filename or "unknown"
    content = await file.read()
    ext = validate_upload(filename, content)
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext or ".tmp") as tmp:
        tmp.write(content)
        path = tmp.name
    return SavedUpload(filename=filename, ext=ext, path=path)


def cleanup_temp_file(path: str | None) -> None:
    if not path:
        return
    try:
        os.unlink(path)
    except OSError:
        pass


def extract_text_and_warnings(file_path: str) -> tuple[str, list[str]]:
    ext = Path(file_path).suffix.lower()
    warnings: list[str] = []
    if ext == ".docx":
        return _extract_docx(file_path, warnings)
    if ext == ".pptx":
        return _extract_pptx(file_path, warnings)
    raise HTTPException(status_code=400, detail=f"不支持的文件格式: {ext}")


def _extract_docx(path: str, warnings: list[str]) -> tuple[str, list[str]]:
    from docx import Document

    doc = Document(path)
    parts: list[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                parts.append(" | ".join(row_texts))

    has_image = any("image" in rel.reltype for rel in doc.part.rels.values())
    if has_image:
        warnings.append("文档中包含图片，暂不支持 OCR 识别图片中的文字，请手动核对")

    return "\n".join(parts), warnings


def _extract_pptx(path: str, warnings: list[str]) -> tuple[str, list[str]]:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    has_image = False

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
            try:
                if "Picture" in type(shape).__name__ or hasattr(shape, "image"):
                    has_image = True
            except Exception:
                pass
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        parts.append(" | ".join(row_texts))

    if has_image:
        warnings.append("演示文稿中包含图片，暂不支持 OCR 识别图片中的文字，请手动核对")

    return "\n".join(parts), warnings


def extract_text_from_file(file_path: str) -> str:
    text, _ = extract_text_and_warnings(file_path)
    return text


def extract_text_or_raise(file_path: str) -> tuple[str, list[str]]:
    try:
        text, warnings = extract_text_and_warnings(file_path)
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"文件提取失败: {exc}") from exc

    if not text.strip():
        raise HTTPException(status_code=400, detail="文档中未提取到任何文本内容")
    return text, warnings


def build_ai_prompt(text_chunk: str) -> str:
    return (
        "You are an exam-question extraction assistant. Convert the document text into strict JSON.\n"
        "Return ONLY a JSON object with this shape:\n"
        "{\n"
        '  "questions": [\n'
        "    {\n"
        '      "type": "single_choice | multiple_choice | true_false | fill_blank | short_answer",\n'
        '      "question": "question text",\n'
        '      "options": {"A": "option A", "B": "option B"},\n'
        '      "answer": "correct answer",\n'
        '      "analysis": "short explanation",\n'
        '      "subject": "subject name",\n'
        '      "chapter": "chapter name",\n'
        '      "difficulty": "easy | normal | hard"\n'
        "    }\n"
        "  ]\n"
        "}\n\n"
        "Rules:\n"
        "1. Use the original language of the document for question text and analysis.\n"
        "2. For choice questions, options must be an object keyed by A/B/C/D.\n"
        "3. For true_false answers, use one of: true, false, yes, no.\n"
        "4. Do not include markdown fences or explanations outside JSON.\n\n"
        f"Document text:\n{text_chunk}"
    )


def validate_question_item(item: dict[str, Any]) -> tuple[dict[str, Any] | None, str | None]:
    q_type = (item.get("type") or "").strip()
    question = (item.get("question") or "").strip()
    answer = (item.get("answer") or "").strip()

    if q_type not in VALID_QUESTION_TYPES:
        return None, f"无效的题目类型 '{q_type}'"
    if not question:
        return None, "题目题干不能为空"
    if not answer:
        return None, "题目答案不能为空"
    if q_type in ("single_choice", "multiple_choice"):
        options = item.get("options")
        if not options or not isinstance(options, dict) or len(options) < 2:
            return None, f"选择题（{q_type}）必须至少提供两个选项"

    return {
        "type": q_type,
        "question": question,
        "options": item.get("options"),
        "answer": normalize_answer(answer, q_type),
        "analysis": (item.get("analysis") or "").strip(),
        "subject": (item.get("subject") or "默认科目").strip(),
        "chapter": (item.get("chapter") or "默认章节").strip(),
        "difficulty": (item.get("difficulty") or "normal").strip(),
    }, None


def question_items_from_parsed_json(parsed: Any) -> list[dict[str, Any]]:
    if isinstance(parsed, list):
        return parsed
    if isinstance(parsed, dict):
        items = (
            parsed.get("questions")
            or parsed.get("items")
            or parsed.get("data")
            or parsed.get("result")
            or []
        )
        if not items and parsed and ("question" in parsed or "type" in parsed):
            return [parsed]
        return items if isinstance(items, list) else []
    return []


def json_candidates_from_text(raw: str) -> list[str]:
    text = (raw or "").strip()
    candidates: list[str] = []
    if not text:
        return candidates

    candidates.append(text)

    for match in re.finditer(r"```(?:json|JSON)?\s*([\s\S]*?)```", text):
        fenced = match.group(1).strip()
        if fenced:
            candidates.append(fenced)

    for open_char, close_char in (("{", "}"), ("[", "]")):
        start = text.find(open_char)
        while start != -1:
            depth = 0
            in_string = False
            escaped = False
            for idx in range(start, len(text)):
                ch = text[idx]
                if in_string:
                    if escaped:
                        escaped = False
                    elif ch == "\\":
                        escaped = True
                    elif ch == '"':
                        in_string = False
                    continue
                if ch == '"':
                    in_string = True
                elif ch == open_char:
                    depth += 1
                elif ch == close_char:
                    depth -= 1
                    if depth == 0:
                        candidates.append(text[start:idx + 1].strip())
                        break
            start = text.find(open_char, start + 1)

    unique: list[str] = []
    seen: set[str] = set()
    for item in candidates:
        if item and item not in seen:
            seen.add(item)
            unique.append(item)
    return unique


def extract_questions_from_ai_response(raw: str) -> tuple[list[dict[str, Any]], list[str]]:
    warnings: list[str] = []
    last_error = ""

    for candidate in json_candidates_from_text(raw):
        try:
            parsed = json_module.loads(candidate)
            if isinstance(parsed, str):
                parsed = json_module.loads(parsed)
        except json_module.JSONDecodeError as exc:
            last_error = str(exc)
            continue

        items = question_items_from_parsed_json(parsed)
        if items:
            if candidate.strip() != (raw or "").strip():
                warnings.append("AI 返回内容包含说明文字，已自动提取其中的 JSON")
            return items, warnings

    if last_error:
        warnings.append(f"AI 返回了非 JSON 格式内容，已忽略此分块（{last_error}）")
    else:
        warnings.append("AI 返回内容中未找到 questions 数组，已忽略此分块")
    return [], warnings


def call_ai_parse_chunk(text_chunk: str, chunk_index: int) -> tuple[list[dict[str, Any]], list[str]]:
    client = OpenAI(
        api_key=OPENAI_API_KEY,
        base_url=OPENAI_BASE_URL,
        timeout=IMPORT_UPSTREAM_TIMEOUT,
    )
    try:
        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": build_ai_prompt(text_chunk)}],
            response_format={"type": "json_object"},
            temperature=0.1,
            max_tokens=3000,
        )
    except APITimeoutError as exc:
        raise HTTPException(status_code=504, detail="AI 调用超时，请稍后重试") from exc
    except TimeoutError as exc:
        raise HTTPException(status_code=504, detail="AI 调用超时，请稍后重试") from exc
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=502, detail=f"AI 调用失败: {exc}") from exc

    raw = response.choices[0].message.content if response.choices else ""
    if not raw:
        return [], [f"第 {chunk_index + 1} 部分 AI 返回了空响应"]

    items, warnings = extract_questions_from_ai_response(raw)
    if warnings:
        warnings = [f"第 {chunk_index + 1} 部分: {warning}" for warning in warnings]
    return items, warnings


def deduplicate_questions(questions: list[dict[str, Any]]) -> list[dict[str, Any]]:
    seen: set[str] = set()
    result: list[dict[str, Any]] = []
    for item in questions:
        key = (item.get("question") or "").strip()[:100]
        if key and key not in seen:
            seen.add(key)
            result.append(item)
    return result


def chunk_document_text(text: str) -> tuple[list[str], int]:
    chunk_start = time.perf_counter()
    paragraphs = [paragraph for paragraph in text.split("\n") if paragraph.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        if len(current) + len(paragraph) + 1 > AI_CHUNK_SIZE:
            if current:
                chunks.append(current)
            current = paragraph
        else:
            current = current + "\n" + paragraph if current else paragraph
    if current:
        chunks.append(current)
    return chunks, elapsed_ms(chunk_start)


def call_ai_parse(text: str) -> tuple[list[dict[str, Any]], list[str], dict[str, Any]]:
    total_start = time.perf_counter()
    timing = {
        "chunk_ms": 0,
        "ai_ms": 0,
        "total_ms": 0,
        "chunks": 0,
        "ai_chunks": [],
    }
    if not OPENAI_API_KEY:
        raise HTTPException(
            status_code=400,
            detail="未配置 OPENAI_API_KEY，请在 .env 文件中设置 OPENAI_API_KEY 以使用 AI 自动导入功能",
        )

    all_items: list[dict[str, Any]] = []
    all_warnings: list[str] = []
    chunks, timing["chunk_ms"] = chunk_document_text(text)

    if len(chunks) > MAX_CHUNKS:
        all_warnings.append(f"文档过长，仅处理前 {MAX_CHUNKS} 部分（共 {len(chunks)} 部分）")
        chunks = chunks[:MAX_CHUNKS]
    timing["chunks"] = len(chunks)

    saw_timeout = False
    saw_invalid_json = False
    other_http_error: HTTPException | None = None

    for index, chunk in enumerate(chunks):
        ai_start = time.perf_counter()
        try:
            items, warnings = call_ai_parse_chunk(chunk, index)
            all_items.extend(items)
            all_warnings.extend(warnings)
            if any("非 JSON" in warning for warning in warnings):
                saw_invalid_json = True
        except HTTPException as exc:
            if exc.status_code == 504:
                saw_timeout = True
            elif exc.status_code not in (400, 502):
                other_http_error = exc
            all_warnings.append(f"第 {index + 1} 部分解析失败: {exc.detail}")
        finally:
            ai_chunk_ms = elapsed_ms(ai_start)
            timing["ai_chunks"].append(ai_chunk_ms)
            timing["ai_ms"] += ai_chunk_ms

    all_items = deduplicate_questions(all_items)

    valid: list[dict[str, Any]] = []
    for index, item in enumerate(all_items):
        validated, error = validate_question_item(item)
        if validated:
            validated["line_number"] = index + 1
            valid.append(validated)
        else:
            all_warnings.append(f"第 {index + 1} 题格式有误: {error}")

    timing["total_ms"] = elapsed_ms(total_start)

    if not valid and other_http_error:
        raise other_http_error
    if not valid and saw_timeout:
        raise HTTPException(status_code=504, detail="AI 调用超时，请稍后重试")
    if not valid and saw_invalid_json:
        raise HTTPException(status_code=400, detail="AI 返回非 JSON，无法解析题目")

    return valid, all_warnings, timing


def ensure_questions_found(questions: list[dict[str, Any]], warnings: list[str]) -> None:
    if questions:
        return
    detail = "未能从文档中解析出任何有效题目"
    if warnings:
        detail += "（" + "；".join(warnings[:3]) + "）"
    raise HTTPException(status_code=400, detail=detail)


def preview_import_from_text(text: str) -> tuple[list[dict[str, Any]], list[str], dict[str, Any]]:
    questions, warnings, timing = call_ai_parse(text)
    ensure_questions_found(questions, warnings)
    return questions, warnings, timing


def validate_imported_questions(
    questions: list[schemas.ImportedQuestion],
) -> tuple[list[schemas.ImportedQuestion], list[str]]:
    errors: list[str] = []
    validated_items: list[schemas.ImportedQuestion] = []
    for index, question in enumerate(questions):
        validated, error = validate_question_item(question.model_dump())
        if error:
            errors.append(f"第 {index + 1} 题: {error}")
        else:
            validated_items.append(schemas.ImportedQuestion(**validated))
    return validated_items, errors


def resolve_target_course(
    db: Session,
    user_id: int,
    *,
    course_id: int = 0,
    course_name: str = "",
    filename: str = "",
):
    try:
        if course_id > 0:
            bank, _ = crud.resolve_course(db, user_id, course_id=course_id)
        elif course_name.strip():
            bank, _ = crud.resolve_course(db, user_id, course_name=course_name.strip())
        else:
            derived = derive_course_name_from_filename(filename)
            bank, _ = crud.resolve_course(db, user_id, course_name=derived)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    return bank


def persist_imported_questions(
    db: Session,
    *,
    user_id: int,
    course_id: int,
    questions: list[schemas.ImportedQuestion | dict[str, Any]],
) -> int:
    imported = 0
    try:
        for question_data in questions:
            data = question_data.model_dump() if hasattr(question_data, "model_dump") else question_data
            question = QuestionModel(
                owner_id=user_id,
                course_id=course_id,
                visibility="private",
                source="import",
                created_at=datetime.now(timezone.utc),
                subject=data["subject"],
                chapter=data["chapter"],
                type=data["type"],
                question=data["question"],
                answer=normalize_answer(data["answer"], data["type"]),
                analysis=data.get("analysis", ""),
                difficulty=data.get("difficulty", "normal"),
            )
            question.set_options_dict(data.get("options"))
            db.add(question)
            imported += 1
        db.commit()
    except Exception as exc:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"导入失败，已回滚: {exc}") from exc
    return imported
