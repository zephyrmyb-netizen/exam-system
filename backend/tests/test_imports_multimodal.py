"""Tests for PPTX and image-backed AI import paths."""

import io
from pathlib import Path
from unittest.mock import MagicMock, patch


def _mock_import_ai_response(question_text: str = "Python Q?"):
    import json

    choice = MagicMock()
    choice.message.content = json.dumps(
        {
            "questions": [
                {
                    "type": "single_choice",
                    "question": question_text,
                    "options": {"A": "language", "B": "fruit"},
                    "answer": "A",
                    "analysis": "mocked",
                }
            ]
        }
    )
    completion = MagicMock()
    completion.choices = [choice]
    client = MagicMock()
    client.chat.completions.create.return_value = completion
    return client


def _make_png_bytes() -> bytes:
    from PIL import Image

    image = Image.new("RGB", (8, 8), color=(255, 255, 255))
    buf = io.BytesIO()
    image.save(buf, format="PNG")
    return buf.getvalue()


def _make_pptx_bytes(text: str = "", image_bytes: bytes | None = None) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if text:
        box = slide.shapes.add_textbox(0, 0, 5000000, 1000000)
        box.text = text
    if image_bytes:
        image_stream = io.BytesIO(image_bytes)
        slide.shapes.add_picture(image_stream, 0, 1000000, width=1000000, height=1000000)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pdf_bytes(text: str) -> bytes:
    # Minimal text PDF with one Helvetica text stream. It keeps the test independent
    # from optional PDF writer libraries while still exercising real PDF parsing.
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    stream = f"BT /F1 14 Tf 72 720 Td ({escaped}) Tj ET".encode("latin-1", errors="ignore")
    objects.append(b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream")

    output = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, obj in enumerate(objects, start=1):
        offsets.append(len(output))
        output.extend(f"{index} 0 obj\n".encode("ascii"))
        output.extend(obj)
        output.extend(b"\nendobj\n")
    xref_offset = len(output)
    output.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    output.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        output.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    output.extend(
        f"trailer\n<< /Root 1 0 R /Size {len(objects) + 1} >>\nstartxref\n{xref_offset}\n%%EOF\n".encode(
            "ascii"
        )
    )
    return bytes(output)


class TestMultimodalImport:
    PREVIEW_URL = "/imports/file/preview"
    FILE_URL = "/imports/file"

    def test_upload_pptx_extracts_slide_text(self, client, auth_headers):
        content = _make_pptx_bytes("1. Python is what? A. language B. fruit Answer: A")
        resp = client.post(
            self.FILE_URL,
            headers=auth_headers,
            files={
                "file": (
                    "python.pptx",
                    content,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )

        assert resp.status_code == 200
        assert "Python" in resp.json()["text"]
        assert "[Slide 1]" in resp.json()["text"]

    def test_upload_pdf_extracts_page_text(self, client, auth_headers):
        content = _make_pdf_bytes("1. PDF question A. yes B. no Answer: A")
        resp = client.post(
            self.FILE_URL,
            headers=auth_headers,
            files={"file": ("pdf-review.pdf", content, "application/pdf")},
        )

        assert resp.status_code == 200
        assert "PDF question" in resp.json()["text"]
        assert "[Page 1]" in resp.json()["text"]

    @patch("backend.routers.imports.OPENAI_API_KEY", "sk-test")
    def test_preview_pptx_with_embedded_image_uses_multimodal_ai(self, client, auth_headers, monkeypatch):
        mock_client = _mock_import_ai_response("Image question")
        monkeypatch.setattr("backend.services.imports_service._build_import_client", lambda: mock_client)

        content = _make_pptx_bytes("Context text", _make_png_bytes())
        resp = client.post(
            self.PREVIEW_URL,
            headers=auth_headers,
            files={"file": ("image-slide.pptx", content, "application/octet-stream")},
        )

        assert resp.status_code == 200
        data = resp.json()
        assert data["questions"][0]["question"] == "Image question"
        assert not any("OCR" in warning and "not support" in warning.lower() for warning in data["warnings"])

        call_kwargs = mock_client.chat.completions.create.call_args.kwargs
        content_parts = call_kwargs["messages"][0]["content"]
        assert isinstance(content_parts, list)
        assert any(part.get("type") == "image_url" for part in content_parts)

    @patch("backend.routers.imports.OPENAI_API_KEY", "sk-test")
    def test_preview_direct_png_uses_image_ai_and_filename_course_name(self, client, auth_headers, monkeypatch):
        mock_client = _mock_import_ai_response("Direct image question")
        monkeypatch.setattr("backend.services.imports_service._build_import_client", lambda: mock_client)

        resp = client.post(
            self.PREVIEW_URL,
            headers=auth_headers,
            files={"file": ("network-quiz.png", _make_png_bytes(), "image/png")},
        )

        assert resp.status_code == 200
        data = resp.json()
        assert data["suggested_course_name"] == "network-quiz"
        assert data["questions"][0]["question"] == "Direct image question"

    def test_upload_legacy_ppt_returns_clear_error(self, client, auth_headers):
        resp = client.post(
            self.FILE_URL,
            headers=auth_headers,
            files={"file": ("legacy.ppt", b"not-a-real-ppt", "application/vnd.ms-powerpoint")},
        )

        assert resp.status_code == 400
        assert ".pptx" in resp.json()["detail"]

    def test_extract_images_from_pptx_limits_image_count(self, tmp_path):
        from backend.imports.file_parser import extract_images_from_file

        image_bytes = _make_png_bytes()
        content = _make_pptx_bytes()
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slide = prs.slides[0]
        for _ in range(13):
            slide.shapes.add_picture(io.BytesIO(image_bytes), 0, 0, width=100000, height=100000)
        path = Path(tmp_path) / "many-images.pptx"
        prs.save(path)

        images, warnings = extract_images_from_file(str(path))

        assert len(images) == 12
        assert any("12" in warning for warning in warnings)
