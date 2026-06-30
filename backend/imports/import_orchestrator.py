import json as json_module
import os
import re
import tempfile
import time
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from fastapi import HTTPException, UploadFile
from openai import APIConnectionError, APIStatusError, OpenAI
from sqlalchemy.orm import Session

from .. import crud, schemas
from ..config import (
    IMPORT_CHUNK_SIZE,
    IMPORT_MAX_CHUNKS,
    IMPORT_UPSTREAM_TIMEOUT,
    OPENAI_API_KEY,
    OPENAI_BASE_URL,
    OPENAI_MODEL,
)
from ..crud import derive_course_name_from_filename
from ..models import Question as QuestionModel
from ..utils import VALID_QUESTION_TYPES, normalize_answer
from .image_extractor import IMAGE_EXTENSIONS, ImagePayload, extract_images_from_file, image_bytes_to_data_url

try:
    from openai import APITimeoutError
except ImportError:  # pragma: no cover
    APITimeoutError = TimeoutError


ALLOWED_EXTENSIONS = {".docx", ".pdf", ".pptx", ".png", ".jpg", ".jpeg", ".webp"}
LEGACY_PPT_EXTENSION = ".ppt"
MAX_FILE_SIZE = 10 * 1024 * 1024
AI_CHUNK_SIZE = max(1000, IMPORT_CHUNK_SIZE)
MAX_CHUNKS = max(1, IMPORT_MAX_CHUNKS)


@dataclass
class SavedUpload:
    filename: str
    ext: str
    path: str


def elapsed_ms(start: float) -> int:
    return max(0, round((time.perf_counter() - start) * 1000))


def build_timing(
    *,
    total_start: float,
    extract_ms: int = 0,
    parse_timing: dict | None = None,
) -> schemas.ImportTiming:
    parse_timing = parse_timing or {}
    return schemas.ImportTiming(
        extract_ms=extract_ms,
        chunk_ms=int(parse_timing.get("chunk_ms") or 0),
        ai_ms=int(parse_timing.get("ai_ms") or 0),
        total_ms=elapsed_ms(total_start),
        chunks=int(parse_timing.get("chunks") or 0),
        ai_chunks=list(parse_timing.get("ai_chunks") or []),
    )


def _ai_override_active() -> bool:
    """Return True when tests have injected a fake OpenAI class."""
    return OpenAI is not __import__("openai").OpenAI


def validate_upload(filename: str, content: bytes) -> str:
    ext = Path(filename or "").suffix.lower()
    if ext == LEGACY_PPT_EXTENSION:
        raise HTTPException(
            status_code=400,
            detail="暂不支持旧版 .ppt，请在 PowerPoint/WPS 中另存为 .pptx 后上传。",
        )
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件格式 '{ext}'，仅支持 .docx、.pdf、.pptx、.png、.jpg、.jpeg、.webp 文件",
        )
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"文件大小超过限制（最大 {MAX_FILE_SIZE // (1024 * 1024)}MB）",
        )
    return ext


def detect_file_kind(filename_or_path: str) -> str:
    ext = Path(filename_or_path or "").suffix.lower()
    if ext == ".docx":
        return "docx"
    if ext == ".pdf":
        return "pdf"
    if ext == ".pptx":
        return "pptx"
    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext == LEGACY_PPT_EXTENSION:
        return "legacy_ppt"
    return "unsupported"


async def save_upload_to_temp(file: UploadFile) -> SavedUpload:
    filename = file.filename or "unknown"
    content = await file.read()
    ext = validate_upload(filename, content)
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext or ".tmp") as tmp:
        tmp.write(content)
        path = tmp.name
    return SavedUpload(filename=filename, ext=ext, path=path)


def cleanup_temp_file(path: str | None) -> None:
    if not path:
        return
    try:
        os.unlink(path)
    except OSError:
        pass


def extract_text_and_warnings(file_path: str) -> tuple[str, list[str]]:
    ext = Path(file_path).suffix.lower()
    warnings: list[str] = []
    if ext == ".docx":
        return _extract_docx(file_path, warnings)
    if ext == ".pdf":
        return _extract_pdf(file_path, warnings)
    if ext == ".pptx":
        return _extract_pptx(file_path, warnings)
    if ext in IMAGE_EXTENSIONS:
        return "", warnings
    if ext == LEGACY_PPT_EXTENSION:
        raise HTTPException(status_code=400, detail="暂不支持旧版 .ppt，请在 PowerPoint/WPS 中另存为 .pptx 后上传。")
    raise HTTPException(status_code=400, detail=f"不支持的文件格式: {ext}")


def _extract_docx(path: str, warnings: list[str]) -> tuple[str, list[str]]:
    from docx import Document

    doc = Document(path)
    parts: list[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                parts.append(" | ".join(row_texts))

    return "\n".join(parts), warnings


def _extract_pdf(path: str, warnings: list[str]) -> tuple[str, list[str]]:
    from pypdf import PdfReader

    try:
        reader = PdfReader(path)
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"PDF 文件读取失败：{exc}") from exc

    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception as exc:
            raise HTTPException(status_code=400, detail="PDF 文件已加密，请先解除密码后再上传。") from exc

    parts: list[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception as exc:
            warnings.append(f"第 {page_index} 页 PDF 文本提取失败：{exc}")
            continue
        if text:
            parts.append(f"[Page {page_index}]\n{text}")

    return "\n\n".join(parts), warnings


def _extract_pptx(path: str, warnings: list[str]) -> tuple[str, list[str]]:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        slide_parts.append(" | ".join(row_texts))
        if slide_parts:
            parts.append(f"[Slide {slide_index}]\n" + "\n".join(slide_parts))

    return "\n".join(parts), warnings


def extract_text_from_file(file_path: str) -> str:
    text, _ = extract_text_and_warnings(file_path)
    return text


def _extract_text_or_raise_legacy(file_path: str) -> tuple[str, list[str]]:
    try:
        text, warnings = extract_text_and_warnings(file_path)
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"文件提取失败: {exc}") from exc

    if not text.strip():
        raise HTTPException(status_code=400, detail="文档中未提取到任何文本内容")
    return text, warnings


def extract_text_or_raise(file_path: str) -> tuple[str, list[str]]:
    try:
        text, warnings = extract_text_and_warnings(file_path)
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"文件提取失败: {exc}") from exc

    if text.strip():
        return text, warnings

    images, image_warnings = extract_images_from_file(file_path)
    warnings.extend(image_warnings)
    ext = Path(file_path).suffix.lower()
    if images:
        return text, warnings
    if ext == ".pptx":
        raise HTTPException(status_code=400, detail="未从 PPT 中识别到文字或图片题目，请检查文件是否为空，或尝试导出为图片后上传。")
    if ext == ".pdf":
        raise HTTPException(status_code=400, detail="未从 PDF 中提取到文字，请确认不是扫描版图片 PDF，或将页面导出为图片后上传。")
    raise HTTPException(status_code=400, detail="文档中未提取到任何文本内容")


def build_ai_prompt(text_chunk: str) -> str:
    return (
        "You are an exam-question extraction assistant. Convert the document text into strict JSON.\n"
        "Return ONLY a JSON object with this shape:\n"
        "{\n"
        '  "questions": [\n'
        "    {\n"
        '      "type": "single_choice | multiple_choice | true_false | fill_blank | short_answer",\n'
        '      "question": "question text",\n'
        '      "options": {"A": "option A", "B": "option B"},\n'
        '      "answer": "correct answer",\n'
        '      "analysis": "short explanation",\n'
        '      "subject": "subject name",\n'
        '      "chapter": "chapter name",\n'
        '      "difficulty": "easy | normal | hard"\n'
        "    }\n"
        "  ]\n"
        "}\n\n"
        "Rules:\n"
        "1. Use the original language of the document for question text and analysis.\n"
        "2. For choice questions, options must be an object keyed by A/B/C/D.\n"
        "3. For true_false answers, use one of: true, false, yes, no.\n"
        "4. Do not include markdown fences or explanations outside JSON.\n\n"
        f"Document text:\n{text_chunk}"
    )


def build_ai_multimodal_prompt(text: str = "") -> str:
    context = text.strip() or "(no extracted text; identify questions from the images)"
    return (
        "You are an exam-question extraction assistant. Extract all exam questions from the text and images.\n"
        "Return ONLY strict JSON, with no markdown or explanation. Shape:\n"
        '{"questions":[{"type":"single_choice | multiple_choice | true_false | fill_blank | short_answer",'
        '"question":"question text","options":{"A":"option A","B":"option B"},'
        '"answer":"correct answer","analysis":"short explanation","subject":"",'
        '"chapter":"","difficulty":"normal"}]}\n'
        "Rules: single_choice answer like A; multiple_choice answer like A,B; true_false answer uses "
        "正确 or 错误; extract every question from every image; skip uncertain decorative text.\n\n"
        f"Extracted document text:\n{context}"
    )


def build_ai_image_text_prompt() -> str:
    return (
        "Identify the exam-question text in these images. Return plain text only, preserving question stems, "
        "options, answers, and analysis when visible. Do not describe decorative elements."
    )


def build_ai_repair_prompt(raw_response: str) -> str:
    return (
        "The previous assistant response was not valid import JSON. "
        "Convert it into strict JSON now.\n"
        "Return ONLY a JSON object with this shape:\n"
        '{"questions":[{"type":"single_choice | multiple_choice | true_false | fill_blank | short_answer",'
        '"question":"question text","options":{"A":"option A","B":"option B"},'
        '"answer":"correct answer","analysis":"short explanation","subject":"subject name",'
        '"chapter":"chapter name","difficulty":"easy | normal | hard"}]}\n'
        'If the text contains no question, return {"questions":[]}.\n\n'
        f"Previous response:\n{raw_response}"
    )


def validate_question_item(item: dict[str, Any]) -> tuple[dict[str, Any] | None, str | None]:
    q_type = (item.get("type") or "").strip()
    question = (item.get("question") or "").strip()
    answer = (item.get("answer") or "").strip()

    if q_type not in VALID_QUESTION_TYPES:
        return None, f"无效的题目类型 '{q_type}'"
    if not question:
        return None, "题目题干不能为空"
    if not answer:
        return None, "题目答案不能为空"
    if q_type in ("single_choice", "multiple_choice"):
        options = item.get("options")
        if not options or not isinstance(options, dict) or len(options) < 2:
            return None, f"选择题（{q_type}）必须至少提供两个选项"

    return {
        "type": q_type,
        "question": question,
        "options": item.get("options"),
        "answer": normalize_answer(answer, q_type),
        "analysis": (item.get("analysis") or "").strip(),
        "subject": (item.get("subject") or "默认科目").strip(),
        "chapter": (item.get("chapter") or "默认章节").strip(),
        "difficulty": (item.get("difficulty") or "normal").strip(),
    }, None


def question_items_from_parsed_json(parsed: Any) -> list[dict[str, Any]]:
    if isinstance(parsed, list):
        return parsed
    if isinstance(parsed, dict):
        items = parsed.get("questions") or parsed.get("items") or parsed.get("data") or parsed.get("result") or []
        if not items and parsed and ("question" in parsed or "type" in parsed):
            return [parsed]
        return items if isinstance(items, list) else []
    return []


def json_candidates_from_text(raw: str) -> list[str]:
    text = (raw or "").strip()
    candidates: list[str] = []
    if not text:
        return candidates

    candidates.append(text)

    for match in re.finditer(r"```(?:json|JSON)?\s*([\s\S]*?)```", text):
        fenced = match.group(1).strip()
        if fenced:
            candidates.append(fenced)

    for open_char, close_char in (("{", "}"), ("[", "]")):
        start = text.find(open_char)
        while start != -1:
            depth = 0
            in_string = False
            escaped = False
            for idx in range(start, len(text)):
                ch = text[idx]
                if in_string:
                    if escaped:
                        escaped = False
                    elif ch == "\\":
                        escaped = True
                    elif ch == '"':
                        in_string = False
                    continue
                if ch == '"':
                    in_string = True
                elif ch == open_char:
                    depth += 1
                elif ch == close_char:
                    depth -= 1
                    if depth == 0:
                        candidates.append(text[start : idx + 1].strip())
                        break
            start = text.find(open_char, start + 1)

    unique: list[str] = []
    seen: set[str] = set()
    for item in candidates:
        if item and item not in seen:
            seen.add(item)
            unique.append(item)
    return unique


def extract_questions_from_ai_response(raw: str) -> tuple[list[dict[str, Any]], list[str]]:
    warnings: list[str] = []
    last_error = ""

    for candidate in json_candidates_from_text(raw):
        try:
            parsed = json_module.loads(candidate)
            if isinstance(parsed, str):
                parsed = json_module.loads(parsed)
        except json_module.JSONDecodeError as exc:
            last_error = str(exc)
            continue

        items = question_items_from_parsed_json(parsed)
        if items:
            if candidate.strip() != (raw or "").strip():
                warnings.append("AI 返回内容包含说明文字，已自动提取其中的 JSON")
            return items, warnings

    if last_error:
        warnings.append(f"AI 返回了非 JSON 格式内容，已忽略此分块（{last_error}）")
    else:
        warnings.append("AI 返回内容中未找到 questions 数组，已忽略此分块")
    return [], warnings


def _build_import_client():
    """Return the OpenAI client to use for import parsing.

    Existing tests inject a fake ``OpenAI`` class via ``sync_ai_settings``;
    when that override is active we honor it (instantiating the injected
    class per-call, as before) so mocks keep working. In production there
    is no override, so we reuse the singleton import client from
    ``ai_client`` to avoid rebuilding the httpx connection pool each call.
    """
    if _ai_override_active():
        return OpenAI(
            api_key=OPENAI_API_KEY,
            base_url=OPENAI_BASE_URL,
            timeout=IMPORT_UPSTREAM_TIMEOUT,
        )
    from ..services.ai_client import get_import_client

    return get_import_client()


def safe_ai_error_detail(exc: Exception | None = None) -> str:
    """Return a user-safe upstream AI error message.

    Provider exceptions can include request ids, upstream internals, or even
    credential fragments. Keep API responses useful without leaking details.
    """
    if isinstance(exc, APIStatusError):
        if exc.status_code in (401, 403):
            return "AI 服务鉴权失败，请检查后端 AI Key 和接口地址配置"
        if exc.status_code == 429:
            return "AI 服务请求过于频繁，请稍后重试"
    if isinstance(exc, APIConnectionError):
        return "AI 服务连接失败，请检查网络或接口地址后重试"
    return "AI 服务暂时不可用，请稍后重试"


def _ensure_openai_key() -> None:
    if not OPENAI_API_KEY:
        raise HTTPException(
            status_code=400,
            detail="未配置 OPENAI_API_KEY，请在 .env 文件中设置 OPENAI_API_KEY 以使用 AI 自动导入功能",
        )


def _multimodal_content(prompt: str, images: list[ImagePayload]) -> list[dict[str, Any]]:
    content: list[dict[str, Any]] = [{"type": "text", "text": prompt}]
    for image in images:
        content.append(
            {
                "type": "image_url",
                "image_url": {"url": image_bytes_to_data_url(image.data, image.mime_type)},
            }
        )
    return content


def _call_chat_completion(content: str | list[dict[str, Any]], *, temperature: float = 0.1):
    client = _build_import_client()
    try:
        return client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": content}],
            response_format={"type": "json_object"},
            temperature=temperature,
            max_tokens=3000,
        )
    except APITimeoutError as exc:
        raise HTTPException(status_code=504, detail="AI 调用超时，请稍后重试") from exc
    except TimeoutError as exc:
        raise HTTPException(status_code=504, detail="AI 调用超时，请稍后重试") from exc
    except HTTPException:
        raise
    except (APIStatusError, APIConnectionError) as exc:
        raise HTTPException(status_code=502, detail=safe_ai_error_detail(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=502, detail=safe_ai_error_detail(exc)) from exc


def _validate_ai_items(items: list[dict[str, Any]], warnings: list[str]) -> list[dict[str, Any]]:
    valid: list[dict[str, Any]] = []
    for index, item in enumerate(deduplicate_questions(items)):
        validated, error = validate_question_item(item)
        if validated:
            validated["line_number"] = index + 1
            valid.append(validated)
        else:
            warnings.append(f"第 {index + 1} 题格式有误: {error}")
    return valid


def call_ai_parse_multimodal(
    text: str,
    images: list[ImagePayload],
) -> tuple[list[dict[str, Any]], list[str], dict[str, Any]]:
    _ensure_openai_key()
    total_start = time.perf_counter()
    ai_start = time.perf_counter()
    response = _call_chat_completion(_multimodal_content(build_ai_multimodal_prompt(text), images))
    ai_ms = elapsed_ms(ai_start)
    raw = response.choices[0].message.content if response.choices else ""
    if not raw:
        timing = {"chunk_ms": 0, "ai_ms": ai_ms, "chunks": 1, "ai_chunks": [ai_ms], "total_ms": elapsed_ms(total_start)}
        return [], ["AI 返回了空响应"], timing

    items, warnings = extract_questions_from_ai_response(raw)
    valid = _validate_ai_items(items, warnings)
    timing = {"chunk_ms": 0, "ai_ms": ai_ms, "chunks": 1, "ai_chunks": [ai_ms], "total_ms": elapsed_ms(total_start)}
    if not valid:
        ensure_questions_found(valid, warnings)
    return valid, warnings, timing


def call_ai_extract_text_from_images(images: list[ImagePayload]) -> tuple[str, list[str], dict[str, Any]]:
    _ensure_openai_key()
    total_start = time.perf_counter()
    ai_start = time.perf_counter()
    client = _build_import_client()
    try:
        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": _multimodal_content(build_ai_image_text_prompt(), images)}],
            temperature=0.1,
            max_tokens=3000,
        )
    except APITimeoutError as exc:
        raise HTTPException(status_code=504, detail="AI 调用超时，请稍后重试") from exc
    except TimeoutError as exc:
        raise HTTPException(status_code=504, detail="AI 调用超时，请稍后重试") from exc
    except HTTPException:
        raise
    except (APIStatusError, APIConnectionError) as exc:
        raise HTTPException(status_code=502, detail=safe_ai_error_detail(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=502, detail=safe_ai_error_detail(exc)) from exc

    ai_ms = elapsed_ms(ai_start)
    text = response.choices[0].message.content if response.choices else ""
    timing = {"chunk_ms": 0, "ai_ms": ai_ms, "chunks": 1, "ai_chunks": [ai_ms], "total_ms": elapsed_ms(total_start)}
    if not text.strip():
        return "", ["图片识别失败：AI 返回了空响应"], timing
    return text.strip(), [], timing


def call_ai_parse_chunk(text_chunk: str, chunk_index: int) -> tuple[list[dict[str, Any]], list[str]]:
    client = _build_import_client()
    try:
        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": build_ai_prompt(text_chunk)}],
            response_format={"type": "json_object"},
            temperature=0.1,
            max_tokens=3000,
        )
    except APITimeoutError as exc:
        raise HTTPException(status_code=504, detail="AI 调用超时，请稍后重试") from exc
    except TimeoutError as exc:
        raise HTTPException(status_code=504, detail="AI 调用超时，请稍后重试") from exc
    except HTTPException:
        raise
    except (APIStatusError, APIConnectionError) as exc:
        raise HTTPException(status_code=502, detail=safe_ai_error_detail(exc)) from exc
    except Exception as exc:
        raise HTTPException(status_code=502, detail=safe_ai_error_detail(exc)) from exc

    raw = response.choices[0].message.content if response.choices else ""
    if not raw:
        return [], [f"第 {chunk_index + 1} 部分 AI 返回了空响应"]

    items, warnings = extract_questions_from_ai_response(raw)
    should_repair = not items and raw.strip() and '"questions"' not in raw and '"question"' not in raw
    if should_repair:
        try:
            repair_response = client.chat.completions.create(
                model=OPENAI_MODEL,
                messages=[{"role": "user", "content": build_ai_repair_prompt(raw)}],
                response_format={"type": "json_object"},
                temperature=0,
                max_tokens=3000,
            )
            repair_raw = repair_response.choices[0].message.content if repair_response.choices else ""
            repaired_items, repair_warnings = extract_questions_from_ai_response(repair_raw)
            if repaired_items:
                items = repaired_items
                warnings = ["AI 返回格式异常，已自动尝试修复为 JSON"] + repair_warnings
        except (APITimeoutError, TimeoutError):
            warnings.append("AI 返回格式异常，自动修复请求超时")
        except Exception:
            warnings.append("AI 返回格式异常，自动修复失败")

    if warnings:
        warnings = [f"第 {chunk_index + 1} 部分: {warning}" for warning in warnings]
    return items, warnings


def deduplicate_questions(questions: list[dict[str, Any]]) -> list[dict[str, Any]]:
    seen: set[str] = set()
    result: list[dict[str, Any]] = []
    for item in questions:
        key = (item.get("question") or "").strip()[:100]
        if key and key not in seen:
            seen.add(key)
            result.append(item)
    return result


def chunk_document_text(text: str) -> tuple[list[str], int]:
    chunk_start = time.perf_counter()
    paragraphs = [paragraph for paragraph in text.split("\n") if paragraph.strip()]
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        if len(current) + len(paragraph) + 1 > AI_CHUNK_SIZE:
            if current:
                chunks.append(current)
            current = paragraph
        else:
            current = current + "\n" + paragraph if current else paragraph
    if current:
        chunks.append(current)
    return chunks, elapsed_ms(chunk_start)


def call_ai_parse(text: str) -> tuple[list[dict[str, Any]], list[str], dict[str, Any]]:
    total_start = time.perf_counter()
    timing = {
        "chunk_ms": 0,
        "ai_ms": 0,
        "total_ms": 0,
        "chunks": 0,
        "ai_chunks": [],
    }
    if not OPENAI_API_KEY:
        raise HTTPException(
            status_code=400,
            detail="未配置 OPENAI_API_KEY，请在 .env 文件中设置 OPENAI_API_KEY 以使用 AI 自动导入功能",
        )

    all_items: list[dict[str, Any]] = []
    all_warnings: list[str] = []
    chunks, timing["chunk_ms"] = chunk_document_text(text)

    if len(chunks) > MAX_CHUNKS:
        all_warnings.append(f"文档过长，仅处理前 {MAX_CHUNKS} 部分（共 {len(chunks)} 部分）")
        chunks = chunks[:MAX_CHUNKS]
    timing["chunks"] = len(chunks)

    saw_timeout = False
    saw_invalid_json = False
    saw_upstream_error = False
    other_http_error: HTTPException | None = None

    for index, chunk in enumerate(chunks):
        ai_start = time.perf_counter()
        try:
            items, warnings = call_ai_parse_chunk(chunk, index)
            all_items.extend(items)
            all_warnings.extend(warnings)
            if any("非 JSON" in warning for warning in warnings):
                saw_invalid_json = True
        except HTTPException as exc:
            if exc.status_code == 504:
                saw_timeout = True
            elif exc.status_code == 502:
                saw_upstream_error = True
            elif exc.status_code != 400:
                other_http_error = exc
            all_warnings.append(f"第 {index + 1} 部分解析失败: {exc.detail}")
        finally:
            ai_chunk_ms = elapsed_ms(ai_start)
            timing["ai_chunks"].append(ai_chunk_ms)
            timing["ai_ms"] += ai_chunk_ms

    all_items = deduplicate_questions(all_items)

    valid: list[dict[str, Any]] = []
    for index, item in enumerate(all_items):
        validated, error = validate_question_item(item)
        if validated:
            validated["line_number"] = index + 1
            valid.append(validated)
        else:
            all_warnings.append(f"第 {index + 1} 题格式有误: {error}")

    timing["total_ms"] = elapsed_ms(total_start)

    if not valid and other_http_error:
        raise other_http_error
    if not valid and saw_timeout:
        raise HTTPException(status_code=504, detail="AI 调用超时，请稍后重试")
    if not valid and saw_upstream_error:
        raise HTTPException(status_code=502, detail=safe_ai_error_detail())
    if not valid and saw_invalid_json:
        raise HTTPException(status_code=400, detail="AI 返回非 JSON，无法解析题目")

    return valid, all_warnings, timing


def ensure_questions_found(questions: list[dict[str, Any]], warnings: list[str]) -> None:
    if questions:
        return
    detail = "未能从文档中解析出任何有效题目"
    if warnings:
        detail += "（" + "；".join(warnings[:3]) + "）"
    raise HTTPException(status_code=400, detail=detail)


def preview_import_from_text(text: str) -> tuple[list[dict[str, Any]], list[str], dict[str, Any]]:
    questions, warnings, timing = call_ai_parse(text)
    ensure_questions_found(questions, warnings)
    return questions, warnings, timing


def preview_import_from_file_content(
    text: str,
    images: list[ImagePayload],
) -> tuple[list[dict[str, Any]], list[str], dict[str, Any]]:
    if images:
        try:
            return call_ai_parse_multimodal(text, images)
        except HTTPException as exc:
            if text.strip():
                questions, warnings, timing = call_ai_parse(text)
                warnings.insert(0, f"图片识别失败：{exc.detail}")
                ensure_questions_found(questions, warnings)
                return questions, warnings, timing
            raise

    questions, warnings, timing = call_ai_parse(text)
    ensure_questions_found(questions, warnings)
    return questions, warnings, timing


def validate_imported_questions(
    questions: list[schemas.ImportedQuestion],
) -> tuple[list[schemas.ImportedQuestion], list[str]]:
    errors: list[str] = []
    validated_items: list[schemas.ImportedQuestion] = []
    for index, question in enumerate(questions):
        validated, error = validate_question_item(question.model_dump())
        if error:
            errors.append(f"第 {index + 1} 题: {error}")
        else:
            validated_items.append(schemas.ImportedQuestion(**validated))
    return validated_items, errors


def resolve_target_course(
    db: Session,
    user_id: int,
    *,
    course_id: int = 0,
    course_name: str = "",
    filename: str = "",
):
    try:
        if course_id > 0:
            bank, _ = crud.resolve_course(db, user_id, course_id=course_id)
        elif course_name.strip():
            bank, _ = crud.resolve_course(db, user_id, course_name=course_name.strip())
        else:
            derived = derive_course_name_from_filename(filename)
            bank, _ = crud.resolve_course(db, user_id, course_name=derived)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    return bank


def persist_imported_questions(
    db: Session,
    *,
    user_id: int,
    course_id: int,
    questions: list[schemas.ImportedQuestion | dict[str, Any]],
) -> int:
    now = datetime.now(UTC)
    models_to_add: list[QuestionModel] = []
    for question_data in questions:
        data = question_data.model_dump() if hasattr(question_data, "model_dump") else question_data
        question = QuestionModel(
            owner_id=user_id,
            course_id=course_id,
            visibility="private",
            source="import",
            created_at=now,
            subject=data["subject"],
            chapter=data["chapter"],
            type=data["type"],
            question=data["question"],
            answer=normalize_answer(data["answer"], data["type"]),
            analysis=data.get("analysis", ""),
            difficulty=data.get("difficulty", "normal"),
        )
        question.set_options_dict(data.get("options"))
        models_to_add.append(question)

    try:
        db.add_all(models_to_add)
        db.commit()
    except Exception as exc:
        db.rollback()
        raise HTTPException(status_code=500, detail=f"导入失败，已回滚: {exc}") from exc
    return len(models_to_add)
