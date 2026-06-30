"""Image extraction and normalization helpers for AI import."""

from __future__ import annotations

import base64
import io
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, UnidentifiedImageError

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
IMAGE_MIME_BY_EXT = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
}
MAX_IMAGE_COUNT = 12
MAX_IMAGE_SIDE = 1600


@dataclass
class ImagePayload:
    data: bytes
    mime_type: str
    source: str = ""


def image_bytes_to_data_url(image_bytes: bytes, mime_type: str) -> str:
    encoded = base64.b64encode(image_bytes).decode("ascii")
    return f"data:{mime_type};base64,{encoded}"


def _mime_from_name(name: str, fallback: str = "image/png") -> str:
    return IMAGE_MIME_BY_EXT.get(Path(name or "").suffix.lower(), fallback)


def normalize_image_for_ai(image_bytes: bytes) -> bytes:
    """Resize and re-encode an image so the multimodal request stays small."""
    try:
        with Image.open(io.BytesIO(image_bytes)) as img:
            image = img.copy()
    except UnidentifiedImageError:
        raise ValueError("无法识别图片文件")

    image.thumbnail((MAX_IMAGE_SIDE, MAX_IMAGE_SIDE))
    output = io.BytesIO()
    if image.mode in {"RGBA", "LA", "P"}:
        image = image.convert("RGBA")
        image.save(output, format="PNG", optimize=True)
    else:
        image = image.convert("RGB")
        image.save(output, format="JPEG", quality=85, optimize=True)
    return output.getvalue()


def _normalize_payload(image_bytes: bytes, mime_type: str, source: str) -> ImagePayload:
    normalized = normalize_image_for_ai(image_bytes)
    # normalize_image_for_ai emits PNG for alpha images and JPEG otherwise.
    try:
        with Image.open(io.BytesIO(normalized)) as img:
            fmt = (img.format or "").upper()
    except UnidentifiedImageError:
        fmt = ""
    normalized_mime = "image/png" if fmt == "PNG" else "image/jpeg"
    return ImagePayload(data=normalized, mime_type=normalized_mime, source=source)


def _limit_images(images: list[ImagePayload]) -> tuple[list[ImagePayload], list[str]]:
    if len(images) <= MAX_IMAGE_COUNT:
        return images, []
    return (
        images[:MAX_IMAGE_COUNT],
        [f"图片数量超过 {MAX_IMAGE_COUNT} 张，已仅处理前 {MAX_IMAGE_COUNT} 张。"],
    )


def extract_direct_image(path: str) -> tuple[list[ImagePayload], list[str]]:
    ext = Path(path).suffix.lower()
    if ext not in IMAGE_EXTENSIONS:
        return [], []
    try:
        payload = _normalize_payload(Path(path).read_bytes(), _mime_from_name(path), Path(path).name)
        return [payload], []
    except Exception as exc:
        return [], [f"图片处理失败：{exc}"]


def extract_images_from_docx(path: str) -> tuple[list[ImagePayload], list[str]]:
    from docx import Document

    doc = Document(path)
    images: list[ImagePayload] = []
    warnings: list[str] = []
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype:
            continue
        try:
            part = rel.target_part
            mime_type = getattr(part, "content_type", "") or _mime_from_name(getattr(part, "partname", ""))
            images.append(_normalize_payload(part.blob, mime_type, str(getattr(part, "partname", ""))))
        except Exception as exc:
            warnings.append(f"图片处理失败：{exc}")
    limited, limit_warnings = _limit_images(images)
    return limited, warnings + limit_warnings


def extract_images_from_pptx(path: str) -> tuple[list[ImagePayload], list[str]]:
    from pptx import Presentation

    prs = Presentation(path)
    images: list[ImagePayload] = []
    warnings: list[str] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not hasattr(shape, "image"):
                continue
            try:
                image = shape.image
                images.append(
                    _normalize_payload(
                        image.blob,
                        getattr(image, "content_type", "") or _mime_from_name(getattr(image, "filename", "")),
                        f"Slide {slide_index}",
                    )
                )
            except Exception as exc:
                warnings.append(f"第 {slide_index} 页图片处理失败：{exc}")
    limited, limit_warnings = _limit_images(images)
    return limited, warnings + limit_warnings


def extract_images_from_file(path: str) -> tuple[list[ImagePayload], list[str]]:
    ext = Path(path).suffix.lower()
    if ext == ".docx":
        return extract_images_from_docx(path)
    if ext == ".pptx":
        return extract_images_from_pptx(path)
    if ext in IMAGE_EXTENSIONS:
        return extract_direct_image(path)
    return [], []
